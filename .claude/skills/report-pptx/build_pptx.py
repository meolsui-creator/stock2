# -*- coding: utf-8 -*-
"""
report-pptx 빌더
입력:  reports/{종목명}.md
출력:  reports/pptx/{종목명}.pptx   (pptx 폴더 없으면 생성)

투자증권 리서치센터 톤:
  - 배경 화이트 / 제목 딥네이비(#1F2A44) / 본문 그레이
  - 한화 오렌지(#F37321)는 표지 상단 바·구분선·키 수치 강조에만 절제 사용
  - 표: 헤더행만 딥네이비+흰 글자, 본문행은 얇은 그레이 괘선
  - 한글 폰트는 '맑은 고딕' 하나로 고정 (글자 깨짐 방지)

가드레일:
  - 모든 수치 옆 (출처/기준일) 표기 — 출처 없는 숫자는 싣지 않음(원문 .md 책임)
  - 매수/매도·목표가 단정 표현 금지 — 본 스크립트는 판단 문구를 생성하지 않음
  - 마지막 슬라이드에 '학습용 분석' 명시

사용법:
  python build_pptx.py "삼성전자"
  python build_pptx.py "삼성전자" --root "C:/path/to/project"
"""
import os
import re
import sys
import argparse
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- 디자인 토큰 -------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x2A, 0x44)   # 딥네이비 - 제목/표 헤더
ORANGE = RGBColor(0xF3, 0x73, 0x21)  # 한화 오렌지 - 상단 바/구분선/키 수치
BODY = RGBColor(0x40, 0x40, 0x40)   # 본문 그레이
RULE = RGBColor(0xD9, 0xD9, 0xD9)   # 표 본문행 얇은 괘선
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"                    # 한글 폰트 고정

# 슬라이드 순서(항상 고정) - (캐논 제목, 매칭 키워드들)
SLIDE_PLAN = [
    ("종목 개요", ["개요", "소개", "기업"]),
    ("재무 요약", ["재무", "실적"]),
    ("가격/추세", ["가격", "추세", "차트", "기술"]),
    ("뉴스·심리", ["뉴스", "심리", "센티", "이슈"]),
    ("리스크", ["리스크", "위험"]),
    ("한 줄 종합", ["종합", "결론", "한 줄", "한줄"]),
]

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


# ---- 폰트 헬퍼 ---------------------------------------------------------------
def _apply_font(run, size, color, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 동아시아 폰트도 동일하게 지정 (글자 깨짐 방지)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)
    cs = rPr.find(qn('a:cs'))
    if cs is None:
        cs = rPr.makeelement(qn('a:cs'), {})
        rPr.append(cs)
    cs.set('typeface', FONT)


def _set_cell_border(cell, edges=("L", "R", "T", "B"), color=RULE, width_pt=0.75):
    """셀 테두리(얇은 그레이 괘선) 설정."""
    tcPr = cell._tc.get_or_add_tcPr()
    tag_map = {"L": "a:lnL", "R": "a:lnR", "T": "a:lnT", "B": "a:lnB"}
    for e in edges:
        tag = tag_map[e]
        ln = tcPr.find(qn(tag))
        if ln is not None:
            tcPr.remove(ln)
        ln = tcPr.makeelement(qn(tag), {
            'w': str(int(width_pt * 12700)), 'cap': 'flat',
            'cmpd': 'sng', 'algn': 'ctr'
        })
        fill = ln.makeelement(qn('a:solidFill'), {})
        clr = fill.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (color[0], color[1], color[2])})
        fill.append(clr)
        ln.append(fill)
        tcPr.append(ln)


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ---- 마크다운 파서 -----------------------------------------------------------
def parse_md(text):
    """제목/작성일/섹션(제목→본문 라인들)을 추출."""
    lines = text.splitlines()
    title = None
    created = None
    sections = []  # [(heading, [lines])]
    cur = None
    for raw in lines:
        line = raw.rstrip()
        if title is None and re.match(r'^#\s+', line):
            title = re.sub(r'^#\s+', '', line).strip()
            continue
        m = re.match(r'^(작성일|기준일|작성)\s*[:：]\s*(.+)$', line.strip())
        if created is None and m:
            created = m.group(2).strip()
            continue
        h = re.match(r'^##\s+(.+)$', line)
        if h:
            cur = [h.group(1).strip(), []]
            sections.append(cur)
            continue
        if cur is not None:
            cur[1].append(line)
    return title, created, sections


def match_section(canon_keywords, sections):
    for heading, body in sections:
        for kw in canon_keywords:
            if kw in heading:
                return body
    return None


def parse_table(body_lines):
    """연속된 마크다운 표 라인을 [[행],...] 로 변환. 없으면 None."""
    rows = []
    for line in body_lines:
        s = line.strip()
        if s.startswith('|') and s.endswith('|'):
            if re.match(r'^\|[\s:\-\|]+\|$', s):  # 구분행 (---)
                continue
            cells = [c.strip() for c in s.strip('|').split('|')]
            rows.append(cells)
    return rows if len(rows) >= 2 else None


def parse_images(body_lines):
    """마크다운 이미지 ![alt](path) 의 경로 목록."""
    paths = []
    for line in body_lines:
        for m in re.finditer(r'!\[[^\]]*\]\(([^)]+)\)', line):
            paths.append(m.group(1).strip())
    return paths


def non_table_lines(body_lines):
    """표·이미지를 제외한 본문(불릿/문장) 라인."""
    out = []
    for line in body_lines:
        s = line.strip()
        if not s:
            continue
        if s.startswith('|') and s.endswith('|'):
            continue
        if re.match(r'^!\[[^\]]*\]\([^)]+\)\s*$', s):  # 이미지 단독 라인
            continue
        s = re.sub(r'!\[[^\]]*\]\([^)]+\)', '', s).strip()  # 줄 내 이미지 제거
        if s:
            out.append(s)
    return out


def add_image(slide, img_path, root, top=Inches(1.5)):
    """이미지를 슬라이드 폭에 맞춰 비율 유지로 배치. 성공 시 차지한 bottom(Emu) 반환, 실패 시 None."""
    path = img_path
    if not os.path.isabs(path):
        path = os.path.join(root, img_path)
    if not os.path.isfile(path):
        return None
    from PIL import Image  # python-pptx 의존성에 포함(Pillow)
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = (16, 9)
    max_w = SLIDE_W - 2 * MARGIN
    max_h = SLIDE_H - top - Inches(0.6)
    ratio = min(max_w / iw, max_h / ih)
    w = int(iw * ratio); h = int(ih * ratio)
    left = int((SLIDE_W - w) / 2)
    slide.shapes.add_picture(path, left, top, width=w, height=h)
    return top + h


# ---- 슬라이드 그리기 ---------------------------------------------------------
def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def draw_title_bar(slide):
    """상단 한화 오렌지 바 (콘텐츠 슬라이드 공통)."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.12))
    _solid(bar, ORANGE)


def draw_slide_title(slide, text):
    tb, tf = _add_textbox(slide, MARGIN, Inches(0.4), SLIDE_W - 2 * MARGIN, Inches(0.8))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    _apply_font(r, 26, NAVY, bold=True)
    # 제목 아래 오렌지 구분선
    line = slide.shapes.add_shape(1, MARGIN, Inches(1.15), Inches(1.4), Pt(3))
    _solid(line, ORANGE)


def add_footer(slide, note="학습용 분석 자료입니다 · 매수/매도·목표가 판단을 제공하지 않습니다"):
    tb, tf = _add_textbox(slide, MARGIN, SLIDE_H - Inches(0.45),
                          SLIDE_W - 2 * MARGIN, Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = note
    _apply_font(r, 9, RGBColor(0x90, 0x90, 0x90))


def build_cover(prs, title, created):
    slide = _blank_slide(prs)
    # 표지 상단 오렌지 바 (두껍게)
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.55))
    _solid(bar, ORANGE)
    # 종목명
    tb, tf = _add_textbox(slide, MARGIN, Inches(2.6), SLIDE_W - 2 * MARGIN,
                          Inches(1.4), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    _apply_font(r, 54, NAVY, bold=True)
    # 부제
    tb2, tf2 = _add_textbox(slide, MARGIN, Inches(4.0), SLIDE_W - 2 * MARGIN, Inches(0.6))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "가치투자 관점 리서치 리포트"
    _apply_font(r2, 20, BODY)
    # 작성일
    tb3, tf3 = _add_textbox(slide, MARGIN, Inches(4.7), SLIDE_W - 2 * MARGIN, Inches(0.5))
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = f"작성일 {created}"
    _apply_font(r3, 14, RGBColor(0x70, 0x70, 0x70))
    add_footer(slide)
    return slide


def add_bullets(slide, lines, top=Inches(1.5)):
    tb, tf = _add_textbox(slide, MARGIN, top, SLIDE_W - 2 * MARGIN,
                          SLIDE_H - top - Inches(0.6))
    first = True
    for s in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        text = re.sub(r'^[-*]\s+', '', s)
        is_bullet = bool(re.match(r'^[-*]\s+', s))
        r = p.add_run()
        r.text = ("• " + text) if is_bullet else text
        _apply_font(r, 16, BODY)


def add_table(slide, rows, top=Inches(1.6)):
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    width = SLIDE_W - 2 * MARGIN
    height = Inches(0.45) * n_rows
    gtbl = slide.shapes.add_table(n_rows, n_cols, MARGIN, top, width, height)
    tbl = gtbl.table
    # 기본 스타일 제거(밴딩 끔) — 직접 그린 괘선만 보이도록
    tbl.first_row = True
    tbl.horz_banding = False
    for ri, row in enumerate(rows):
        is_header = (ri == 0)
        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            val = row[ci] if ci < len(row) else ""
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # 배경
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if is_header else WHITE
            # 텍스트
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (is_header or ci > 0) else PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            _apply_font(r, 12, WHITE if is_header else BODY, bold=is_header)
            # 괘선: 헤더는 네이비라 괘선 생략, 본문행만 얇은 그레이
            if not is_header:
                _set_cell_border(cell, edges=("T", "B", "L", "R"), color=RULE, width_pt=0.5)


def build_content(prs, title, body_lines, kind, root="."):
    slide = _blank_slide(prs)
    draw_title_bar(slide)
    draw_slide_title(slide, title)
    if body_lines is None:
        add_bullets(slide, ["(원문 .md에 해당 섹션이 없어 비워둡니다. 출처 있는 내용만 채웁니다.)"])
        add_footer(slide)
        return slide

    table = parse_table(body_lines)
    texts = non_table_lines(body_lines)
    images = parse_images(body_lines)

    if kind == "summary":  # 한 줄 종합 — 강조 박스
        joined = " ".join(texts) if texts else "(종합 의견이 .md에 없습니다.)"
        tb, tf = _add_textbox(slide, MARGIN, Inches(2.6),
                              SLIDE_W - 2 * MARGIN, Inches(2.0), anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = joined
        _apply_font(r, 22, NAVY, bold=True)
        add_footer(slide)
        return slide

    top = Inches(1.5)
    if images:
        # 이미지(차트) 우선 배치 후, 남는 공간에 설명 텍스트
        bottom = add_image(slide, images[0], root, top=top)
        if bottom is None:
            add_bullets(slide, texts + ["(차트 이미지를 찾지 못했습니다: %s)" % images[0]], top=top)
        elif texts:
            add_bullets(slide, texts, top=bottom + Inches(0.15))
    elif table:
        add_table(slide, table, top=top)
        # 표 아래 부가 설명(출처/기준일 등)
        extra = [t for t in texts]
        if extra:
            below = Inches(1.6) + Inches(0.45) * len(table) + Inches(0.2)
            add_bullets(slide, extra, top=below)
    else:
        add_bullets(slide, texts if texts else ["(내용이 비어 있습니다.)"], top=top)
    add_footer(slide)
    return slide


# ---- 메인 -------------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("name", help="종목명 (reports/{종목명}.md)")
    ap.add_argument("--root", default=None, help="프로젝트 루트(기본: 현재 폴더)")
    args = ap.parse_args()

    root = args.root or os.getcwd()
    md_path = os.path.join(root, "reports", f"{args.name}.md")
    if not os.path.isfile(md_path):
        print(f"[오류] 입력 파일이 없습니다: {md_path}", file=sys.stderr)
        sys.exit(1)

    with open(md_path, encoding="utf-8") as f:
        text = f.read()

    title, created, sections = parse_md(text)
    title = title or args.name
    created = created or date.today().isoformat()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1) 표지
    build_cover(prs, title, created)
    # 2~7) 고정 순서 콘텐츠
    for canon, keywords in SLIDE_PLAN:
        body = match_section(keywords, sections)
        kind = "summary" if canon == "한 줄 종합" else "normal"
        build_content(prs, canon, body, kind, root=root)

    out_dir = os.path.join(root, "reports", "pptx")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"{args.name}.pptx")
    prs.save(out_path)
    print(f"[완료] {out_path}  (슬라이드 {len(prs.slides._sldIdLst)}장)")


if __name__ == "__main__":
    main()
